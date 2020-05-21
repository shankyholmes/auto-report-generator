from datetime import date
import os
import numpy as np
from pptx import Presentation
import Plots as plots
from fuzzywuzzy import process
from scipy import stats
import pandas as pd
# import seaborn as sns

class reportGen():
    
    def __init__(self, df, treatment_type = 'IQR', savelocation = os.path.join(os.path.dirname(__file__), 'Report.pptx')):
        self.TRUNCATION = 0.05
        self.df_truncated = self.truncated(df, self.TRUNCATION)
        self.df = self.outlierTreatmentIQR(df, type = treatment_type)
        self.dirname = os.path.dirname(__file__)
        self.savelocation = savelocation
    
    
    def truncated(self, df, percent_truncation = 0.02):
    
        # TRIM DATA AT BOTH ENDS
        no_rows = int(len(df) * percent_truncation)
        df = df.tail(-no_rows)
        df = df.head(-no_rows)
        return df
    
    def outlierTreatmentIQR(self, df, type = 'IQR'):
        
        # IQR BASED OUTLIER TREATMENT
        if type == 'IQR':
            q1 = df.quantile(0.25)
            q3 = df.quantile(0.75)
            IQR = q3 - q1
            df_out = df[~((df < (q1 - 1.5 * IQR)) | (df > (q3 + 1.5 * IQR))).any(axis = 1)]
            return df_out
        
        # ZSCORE BASED OUTLIER TREATMENT
        if type == 'Zscore':
            z = np.abs(stats.zscore(df))
            threshold = 3
            df_out = df[(z < threshold).all(axis = 1)]
            return df_out
    
    def pptGenerator(self):
        df1 = self.df
        df_truncated = self.df_truncated
        template = os.path.join(self.dirname, 'template.pptx')
        
        
        # DEFINE CONSTANTS
        TITLE_AND_SUBTITLE = 0
        CONTENT_SLIDE_1 = 1
        CONTENT_SLIDE_2 = 2
        
        # INITIALIZE PRESENTATION
        prs = Presentation(template)
        
        # LAYOUTS FOR SLIDES
        title_slide_layout = prs.slide_layouts[TITLE_AND_SUBTITLE]
        content_slide_layout_1 = prs.slide_layouts[CONTENT_SLIDE_1]
        content_slide_layout_2 = prs.slide_layouts[CONTENT_SLIDE_2]
        
        # ADD TITLE AND SUBTITLE SLIDE
        slide = prs.slides.add_slide(title_slide_layout)
        slide.placeholders[0].text = 'INCA - recording Report'
        slide.placeholders[1].text = 'Generated on {:%m-%d-%Y}'.format(date.today())
        
        # ADD CONTENT

        # ADD JOINTPLOT FOR QSET
        slide = prs.slides.add_slide(content_slide_layout_2)
        
        # extracting label names
        QSET = process.extractOne('InjCtl_qSetUnBal', df1.columns)[0]
        SPEED = process.extractOne('Epm_nEng', df1.columns)[0]
        
        # top left
        plots.distplot(df1[SPEED])
        slide.placeholders[22].insert_picture('distplot.png')
        slide.placeholders[16].text = f'{SPEED} Distribution plot'
        
        # top right
        plots.jointplot(df1[SPEED], df1[QSET])
        slide.placeholders[23].insert_picture('jointplot.png')
        slide.placeholders[17].text = f'Zone mapping of the data'
        
        # bottom left
        # scatter for qset and speed
        plots.scatterplot(df1[SPEED], df1[QSET])
        slide.placeholders[24].insert_picture('scatterplot.png')
        slide.placeholders[20].text = f'{QSET} and {SPEED} scatter'
        
        # #ADD OTHER PLOTS
        for col in df1.columns:
            slide = prs.slides.add_slide(content_slide_layout_1)

            plots.distplot(df1[col])
            plots.scatterplot(df1[SPEED], df1[col])
            plots.lineplot(df_truncated.index.astype(dtype = 'float64'), df_truncated[col])
            
            # top left
            slide.placeholders[22].insert_picture('scatterplot.png')
            slide.placeholders[16].text = f'Scatter plot for {col}'
    
            # top right
    
            slide.placeholders[23].insert_picture('distplot.png')
            slide.placeholders[17].text = f'Distribution plot for {col}'
    
            # bottom left
            slide.placeholders[24].insert_picture('lineplot.png')
            slide.placeholders[20].text = f'y/t Plot for {col} at {self.TRUNCATION} truncation'
    

        
        
        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        
        prs.save(self.savelocation)
        # self.PPTtoPDF('Report.pptx', self.savelocation)
        
    # def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    #     powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
    #     powerpoint.Visible = 1
    #
    #     if outputFileName[-3:] != 'pdf':
    #         outputFileName = outputFileName + ".pdf"
    #     deck = powerpoint.Presentations.Open(inputFileName)
    #     deck.SaveAs(outputFileName, formatType)  # formatType = 32 for ppt to pdf
    #     deck.Close()
    #     powerpoint.Quit()

            
#
#
# df = pd.read_csv('D:\Designing\Programming\PyCharm\MachineLearning\sortedList.csv', index_col = 0)
#
# ppt = reportGen(df, 'IQR').pptGenerator()